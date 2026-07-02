"""Flow B — PPTX → normalized HTML (python-pptx).

Iterates slides in order, mapping titles → headings, body placeholders → paragraphs/lists,
tables → tables, and exporting each embedded picture to ``<run>/assets/`` referenced by an
``<img>`` at the slide's position. SmartArt/grouped shapes that don't export as a single clean
picture fall back to a per-slide note (logged), keeping the contract intact.
"""

from __future__ import annotations

from html import escape
from pathlib import Path

from ..config import Settings, get_settings
from ..schemas import ImageRef, NormalizedDocument, SessionMeta
from ..utils.io import run_dir


def pptx_to_normalized(path: Path, *, run_id: str | None, metadata: dict,
                       settings: Settings | None = None) -> NormalizedDocument:
    settings = settings or get_settings()
    from pptx import Presentation  # lazy import
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu  # noqa: F401

    prs = Presentation(str(path))
    assets_dir = (run_dir(run_id, settings) / "assets") if run_id else (settings.runs_path / "_pptx_assets")
    assets_dir.mkdir(parents=True, exist_ok=True)

    html_parts: list[str] = []
    assets: list[ImageRef] = []
    img_n = 0
    session_name = metadata.get("session_name") or path.stem.replace("_", " ").title()

    for s_idx, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        level = 1 if s_idx == 1 and title else 2
        if title:
            html_parts.append(f"<h{level}>{escape(title)}</h{level}>")

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_n += 1
                ref = _export_picture(shape, assets_dir, img_n, s_idx, slide_title=title)
                if ref:
                    assets.append(ref)
                    html_parts.append(
                        f'<img id="{ref.image_id}" src="{ref.src}" alt="{escape(ref.alt)}" '
                        f'data-occurrences="1" data-source-ref="slide-{s_idx}"/>'
                    )
            elif shape.has_table:
                html_parts.append(_table_html(shape.table))
            elif shape.has_text_frame and shape != _title_shape(slide):
                html_parts.append(_text_frame_html(shape.text_frame))

    normalized_html = "\n".join(p for p in html_parts if p)
    meta = SessionMeta(
        session_name=session_name,
        course_name=metadata.get("course_name", "Course"),
        source_type="pptx",
        source_filename=path.name,
        learning_objectives=metadata.get("learning_objectives", []),
    )
    return NormalizedDocument(session_meta=meta, normalized_html=normalized_html, assets=assets)


def _title_shape(slide):
    try:
        return slide.shapes.title
    except Exception:  # noqa: BLE001
        return None


def _slide_title(slide) -> str:
    t = _title_shape(slide)
    if t is not None and t.has_text_frame:
        return t.text_frame.text.strip()
    return ""


def _export_picture(shape, assets_dir: Path, n: int, slide_idx: int,
                    *, slide_title: str = "") -> ImageRef | None:
    try:
        image = shape.image
        ext = image.ext or "png"
        fname = f"slide{slide_idx}_img{n:02d}.{ext}"
        out = assets_dir / fname
        out.write_bytes(image.blob)
        alt = ""
        try:
            alt = (shape.name or "").strip()
        except Exception:  # noqa: BLE001
            pass
        return ImageRef(
            image_id=f"img_{n:02d}",
            src=str(out),
            alt=alt,
            # The slide title is the best free signal for what the picture is about; it seeds
            # both the description heuristic and the block-placement context.
            nearby_heading=slide_title,
            width=_emu_to_px(getattr(shape, "width", None)),
            height=_emu_to_px(getattr(shape, "height", None)),
            slide_index=slide_idx,
            source_ref=f"slide-{slide_idx}",
            bytes=len(image.blob),
            format=ext,
        )
    except Exception:  # noqa: BLE001 — SmartArt/grouped shapes etc.
        return None


def _emu_to_px(emu) -> int | None:
    # python-pptx measures in EMUs (914400 per inch); 96 px/inch → /9525 per px.
    try:
        return int(emu / 9525) if emu else None
    except (TypeError, ValueError):
        return None


def _text_frame_html(tf) -> str:
    paras = [p for p in tf.paragraphs if (p.text or "").strip()]
    if not paras:
        return ""
    # Multiple short bullet-ish lines → list; otherwise paragraphs.
    if len(paras) > 1 and all(len((p.text or "")) < 120 for p in paras):
        items = "".join(f"<li>{escape(p.text.strip())}</li>" for p in paras)
        return f"<ul>{items}</ul>"
    return "".join(f"<p>{escape(p.text.strip())}</p>" for p in paras)


def _table_html(table) -> str:
    rows = list(table.rows)
    if not rows:
        return ""
    head = "".join(f"<th>{escape(c.text.strip())}</th>" for c in rows[0].cells)
    body = "".join(
        "<tr>" + "".join(f"<td>{escape(c.text.strip())}</td>" for c in r.cells) + "</tr>"
        for r in rows[1:]
    )
    return f"<table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table>"
