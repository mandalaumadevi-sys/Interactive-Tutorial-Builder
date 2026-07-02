"""PPTX ingestion (Flow B) with an embedded image, end-to-end on the mock LLM."""

import struct
import uuid
import zlib

import pytest

from tutorial_builder.graph import build_graph, initial_state
from tutorial_builder.ingest import ingest


def _make_png(path, w=16, h=16, rgb=(59, 130, 246)):
    raw = b"".join(b"\x00" + bytes(rgb) * w for _ in range(h))

    def chunk(t, d):
        c = t + d
        return struct.pack(">I", len(d)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    path.write_bytes(b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
                     + chunk(b"IDAT", zlib.compress(raw, 9)) + chunk(b"IEND", b""))


def _make_deck(path, png):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    blank = prs.slide_layouts[5]
    slides = [
        ("Software Development Models", "Why teams need different development models in practice today."),
        ("Waterfall Model Lifecycle", "A sequential lifecycle: requirements, design, build, test, deploy."),
        ("Agile Model", "Iterative delivery in short cycles with continuous feedback loops here."),
        ("V-Model Testing", "Each development phase is paired with a matching testing phase here."),
    ]
    for i, (title, body) in enumerate(slides):
        s = prs.slides.add_slide(blank)
        s.shapes.title.text = title
        tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
        tb.text_frame.text = body
        if i == 1:
            s.shapes.add_picture(str(png), Inches(1), Inches(3), Inches(4), Inches(3))
    prs.save(str(path))


def test_pptx_ingest_extracts_and_describes_image(tmp_path):
    png = tmp_path / "diagram.png"
    _make_png(png)
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, png)

    doc = ingest(str(deck), "pptx", run_id=None,
                 metadata={"course_name": "C", "session_name": "SDLC Deck"})
    assert len(doc.assets) == 1
    img = doc.assets[0]
    assert img.description                       # auto-described
    assert img.nearby_heading == "Waterfall Model Lifecycle"
    assert img.animation_worthy is True          # diagram on a "lifecycle/model" slide


def test_pptx_full_pipeline_mock(tmp_path):
    png = tmp_path / "diagram.png"
    _make_png(png)
    deck = tmp_path / "deck.pptx"
    _make_deck(deck, png)

    graph = build_graph()
    cfg = {"configurable": {"thread_id": uuid.uuid4().hex[:8]}}
    meta = {"course_name": "C", "session_name": "SDLC Deck"}
    graph.invoke(initial_state("p1", str(deck), "pptx", meta), cfg)

    accept = {
        "human_block_review": {"blocks_accepted": True},
        "human_content_review": {"content_accepted": True},
        "human_animation_review": {"animation_accepted": True},
        "human_mcq_review": {"mcq_accepted": True},
        "human_assessment_review": {"assessment_accepted": True},
    }
    for _ in range(12):
        nxt = tuple(graph.get_state(cfg).next or ())
        if not nxt:
            break
        for gate, patch in accept.items():
            if gate in nxt:
                graph.update_state(cfg, patch)
                break
        graph.invoke(None, cfg)

    vals = graph.get_state(cfg).values
    assert vals.get("status") == "completed"
    assert "QUIZ_DATA" in (vals.get("final_html") or "")
    # the waterfall diagram became an inline animation
    assert sum(len(b.get("animations", [])) for b in vals.get("built_blocks_list", [])) >= 1
